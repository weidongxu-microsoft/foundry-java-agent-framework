"""Generate demo.pptx — 7-slide code-walk comparison deck (wire -> stream -> flow -> memory)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

BG = RGBColor(0x1E, 0x1E, 0x2E)
FG = RGBColor(0xE8, 0xE8, 0xF0)
MUTED = RGBColor(0xA8, 0xA8, 0xC0)
GREEN = RGBColor(0x6C, 0xCF, 0x8E)   # with framework
AMBER = RGBColor(0xE0, 0xB0, 0x5A)   # without framework
BLUE = RGBColor(0x7A, 0xA2, 0xF7)    # mode: photo-process (slide 2 only)
PURPLE = RGBColor(0xC0, 0x8C, 0xF0)  # mode: chat (slide 2 only)
CODEBG_L = RGBColor(0x18, 0x2A, 0x20)
CODEBG_R = RGBColor(0x2A, 0x24, 0x18)
ACCENT = RGBColor(0x7A, 0xA2, 0xF7)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    return tb, tb.text_frame


def set_run(r, text, size, color, bold=False, mono=False):
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = "Consolas" if mono else "Segoe UI"


def title(slide, text, sub=None):
    _, tf = box(slide, 0.6, 0.32, 12.1, 1.0)
    set_run(tf.paragraphs[0].add_run(), text, 29, FG, bold=True)
    if sub:
        p = tf.add_paragraph()
        set_run(p.add_run(), sub, 15, MUTED)


def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def code_panel(slide, l, t, w, h, header, hcolor, lines, cbg, count=None, note=None):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(l), Inches(t), Inches(w), Inches(h))
    fill(card, cbg)
    _, hf = box(slide, l + 0.15, t + 0.06, w - 0.3, 0.95)
    set_run(hf.paragraphs[0].add_run(), header, 11.5, hcolor, bold=True)
    if count:  # enlarged, loud line-count for the hand-written side
        p = hf.add_paragraph()
        set_run(p.add_run(), count, 21, hcolor, bold=True)
    if note is not None:  # framework writes nothing here
        _, nf = box(slide, l + 0.3, t + h / 2 - 0.55, w - 0.6, 1.3)
        nf.paragraphs[0].alignment = PP_ALIGN.CENTER
        set_run(nf.paragraphs[0].add_run(), "No code needed", 24, hcolor, bold=True)
        p = nf.add_paragraph(); p.alignment = PP_ALIGN.CENTER
        set_run(p.add_run(), note, 13, MUTED)
        return
    top = t + (1.02 if count else 0.58)
    _, cf = box(slide, l + 0.2, top, w - 0.4, h - (top - t) - 0.08)
    for i, ln in enumerate(lines):
        p = cf.paragraphs[0] if i == 0 else cf.add_paragraph()
        p.space_after = Pt(0)
        set_run(p.add_run(), ln if ln else " ", 10.5, FG, mono=True)


def bullets(slide, l, t, w, h, items, size=16):
    _, tf = box(slide, l, t, w, h)
    for i, (txt, color) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        set_run(p.add_run(), txt, size, color)


# ---- Slide 1 — Title ----
s = prs.slides.add_slide(BLANK); bg(s)
_, tf = box(s, 0.8, 2.3, 11.7, 2.6)
set_run(tf.paragraphs[0].add_run(), "One hosted agent, two ways", 44, FG, bold=True)
p = tf.add_paragraph(); set_run(p.add_run(), "a code walk", 28, ACCENT, bold=True)
p = tf.add_paragraph(); p.space_before = Pt(18)
set_run(p.add_run(),
        "Same photography-assistant agent, same Foundry /responses wire contract.", 18, MUTED)
p = tf.add_paragraph()
set_run(p.add_run(), "One slice: the middleware + hosting seam.   ", 18, MUTED)
set_run(p.add_run(), "left = with framework", 18, GREEN, bold=True)
set_run(p.add_run(), "    right = hand-written", 18, AMBER, bold=True)

# ---- Slide 2 — Demo app (shared); modes use blue/purple, NOT green/amber ----
s = prs.slides.add_slide(BLANK); bg(s)
title(s, "The demo app  (shared \u2014 identical on both sides)",
      "A photography-assistant hosted agent. Its middleware routes every turn:")
bullets(s, 0.8, 2.0, 11.7, 3.4, [
    ("\u25a0  Photo attached  \u2192  PHOTO-PROCESS mode", BLUE),
    ("      ask the vision model \u201cwhat's the best way to process this photo for impact?\u201d,", FG),
    ("      get a strict-JSON plan and apply it with the pure-JDK PhotoProcessor.", FG),
    (" ", FG),
    ("\u25a0  No photo  \u2192  CHAT mode", PURPLE),
    ("      a normal streaming turn against the assistant persona.", FG),
], size=17)
bullets(s, 0.8, 5.5, 11.7, 1.4, [
    ("App logic is IDENTICAL on both sides:  PhotoProcessMiddleware \u2248 PhotoProcessWorkflow.", MUTED),
    ("Everything that follows is what each side must write AROUND it.", MUTED),
    ("(Blue/purple = the two modes \u2014 not the with/without split, which is green/amber.)", MUTED),
], size=14)

# ---- Slide 3 — Wire contract (REST + SSE) ----
s = prs.slides.add_slide(BLANK); bg(s)
title(s, "The wire contract  (REST + SSE)",
      "/responses, /readiness, named SSE, response objects: inherited vs hand-written.")
code_panel(s, 0.5, 1.65, 6.1, 4.0,
           "WITH framework \u00b7 inherited components", GREEN, None, CODEBG_L,
           note="/responses, /readiness, named SSE and the whole Response envelope "
                "are inherited beans \u2014 you register them, you implement none of it.")
code_panel(s, 6.8, 1.65, 6.1, 4.0,
           "WITHOUT \u00b7 writeSse + writeFrame + ResponsesJson 154 + Attachment 10", AMBER, [
    "// EVERY frame NAMED; end the stream on",
    "// response.completed, never data:[DONE].",
    "http.setContentType(TEXT_EVENT_STREAM);",
    "http.setHeader(\"X-Accel-Buffering\",\"no\");",
    "PrintWriter w = http.getWriter();",
    "created.put(\"type\",\"response.created\");",
    "created.set(\"response\", responseObject(..));",
    "writeFrame(w,\"response.created\",   created);",
    "delta.put(\"type\",\"response.output_text.delta\");",
    "delta.put(\"delta\", assistantText);",
    "writeFrame(w,\"response.output_text.delta\",delta);",
    "completed.put(\"type\",\"response.completed\");",
    "writeFrame(w,\"response.completed\", completed);",
    "// + ResponsesJson.responseObject(): the whole",
    "//   Response envelope \u2014 item ids, usage, status.",
], CODEBG_R, count="\u2248 220 lines")
bullets(s, 0.5, 5.85, 12.3, 1.0,
        [("The REST/SSE contract, /readiness, and the response object all only exist on the right.",
          MUTED)], size=15)

# ---- Slide 4 — Streaming chat (relay) ----
s = prs.slides.add_slide(BLANK); bg(s)
title(s, "Streaming chat  (the relay)",
      "The model already emits named SSE frames. Who relays them?")
code_panel(s, 0.5, 1.7, 6.1, 4.2,
           "WITH framework \u00b7 the entire streaming path \u2014 1 line", GREEN, [
    "return next.invoke(context);",
    " ",
    "// ChatClient relays the model's own",
    "// named SSE frames internally \u2014 you",
    "// write nothing else on the stream path.",
], CODEBG_L)
code_panel(s, 6.8, 1.7, 6.1, 4.2,
           "WITHOUT \u00b7 openChatStream + relayChat + captureDelta", AMBER, [
    "Stream<String> up =",
    "    model.openChatStream(instr, userText);",
    "http.setContentType(TEXT_EVENT_STREAM);",
    "http.setHeader(\"X-Accel-Buffering\",\"no\");",
    "PrintWriter w = http.getWriter();",
    "StringBuilder asst = new StringBuilder();",
    "try (up) {",
    "  up.forEach(line -> {",
    "    if (\"data: [DONE]\".equals(line.trim()))",
    "        return;        // contract forbids it",
    "    captureDelta(line, asst);  // tap for memory",
    "    w.write(line); w.write(\"\\n\");",
    "    if (line.isEmpty()) w.flush();  // frame end",
    "  });",
    "}",
    "memory.remember(userText, asst.toString(), scope);",
], CODEBG_R, count="\u2248 60 lines")
bullets(s, 0.5, 6.15, 12.3, 0.9,
        [("Framework side writes nothing; without it, you own the verbatim streaming relay by hand.",
          MUTED)], size=15)

# ---- Slide 5 — Routing / middleware flow ----
s = prs.slides.add_slide(BLANK); bg(s)
title(s, "The routing / middleware flow",
      "Both decide the same thing: photo \u2192 photo-process, else hand the turn to the model.")
code_panel(s, 0.5, 1.7, 6.1, 4.2,
           "WITH framework \u00b7 PhotoProcessMiddleware.invoke \u2014 14 lines (full)", GREEN, [
    "DataContent photo = findPhoto(context);",
    "if (photo == null) {",
    "    return next.invoke(context);  // no photo",
    "}",
    "return CompletableFuture.supplyAsync(",
    "  () -> AgentResponse.builder()",
    "     .message(cropPhoto(photo))",
    "     .finishReason(FinishReason.STOP)",
    "     .build(), executor);",
], CODEBG_L)
code_panel(s, 6.8, 1.7, 6.1, 4.2,
           "WITHOUT \u00b7 createResponse \u2014 parse + route by hand", AMBER, [
    "JsonNode in = tree(body.get(\"input\"));",
    "String text  = ResponsesJson.extractText(in);",
    "Attachment photo = ResponsesJson.extractImage(in);",
    "boolean stream = asBoolean(body.get(\"stream\"))",
    "    || accept.contains(TEXT_EVENT_STREAM);",
    "if (photo != null) {              // photo-process",
    "    String out = workflow.run(text, photo);",
    "    if (stream) writeSse(resp, out);",
    "    else        writeJson(resp, out);",
    "    return;",
    "}",
    "// no photo -> chat mode (memory-augmented):",
    "String scope = resolveScope(isoKey, text);",
    "String instr = computeInstructions(text, scope);",
    "if (stream) relayChat(resp, instr, text, scope);",
    "else writeJson(resp, model.completeText(instr,text));",
], CODEBG_R, count="\u2248 40 lines")
bullets(s, 0.5, 6.1, 12.3, 0.9,
        [("Identical decision \u2014 the difference is everything each side writes around it.",
          MUTED)], size=15)

# ---- Slide 6 — Add memory (second worked example) ----
s = prs.slides.add_slide(BLANK); bg(s)
title(s, "Now add memory  (a second worked example)",
      "Durable per-user memory: recall before the turn, remember after.")
code_panel(s, 0.5, 1.6, 6.1, 3.0,
           "WITH framework \u00b7 attach a provider \u2014 1 line", GREEN, [
    "ChatClientAgent.builder(chatClient)",
    "  .aiContextProvider(memoryProvider)",
    "  // ^ recall-before + store-after,",
    "  //   all handled by the framework",
    "  ...",
], CODEBG_L)
code_panel(s, 6.8, 1.6, 6.1, 3.0,
           "WITHOUT \u00b7 MemoryService 241 + ~99 wiring", AMBER, [
    "scope = resolveScope(isoKey, rawText);",
    "  // low-level API resolves no identity",
    "instr = computeInstructions(text, scope);",
    "  // recall + splice into instructions",
    "memory.remember(text, assistant, scope);",
    "  // async extract, secret-guarded;",
    "  // streaming must captureDelta() too",
], CODEBG_R, count="\u2248 340 lines")
bullets(s, 0.5, 4.75, 12.3, 1.3, [
    ("MemoryService still owns the gotchas:  explicit type:\"message\" (data-plane parser)  \u00b7  "
     "secret guard (no redaction)  \u00b7  pleasantry filter (extractor bloat)  \u00b7  recall char-budget.",
     MUTED),
], size=14)
_, tf = box(s, 0.5, 6.15, 12.3, 0.9)
set_run(tf.paragraphs[0].add_run(), "1 line  vs  ~340 lines", 22, ACCENT, bold=True)
set_run(tf.paragraphs[0].add_run(), "   \u2014 and every gotcha is on the right.", 16, MUTED)

# ---- Slide 7 — Just middleware / full surface ----
s = prs.slides.add_slide(BLANK); bg(s)
title(s, "\u2026and this was just middleware",
      "You saw one seam. The framework carries the rest:")
_, tf = box(s, 0.8, 1.6, 11.9, 0.6)
set_run(tf.paragraphs[0].add_run(), "This one demo:  ~340 vs ~880 lines", 20, ACCENT, bold=True)
set_run(tf.paragraphs[0].add_run(),
        "   \u2014 shared app logic aside, the framework removes ~540 lines of plumbing.", 15, MUTED)
bullets(s, 0.8, 2.35, 11.9, 4.4, [
    ("\u2022  12  context / memory providers  \u2014  Foundry memory, Agent Skills, file memory, "
     "compaction, todo, agent-modes, shell, file-access, background agents, code-act, injection.", FG),
    ("\u2022   3  chat providers  \u2014  OpenAI, Foundry (Azure AI), LangChain4j.", FG),
    ("\u2022   9  tool types  \u2014  function, web search, code interpreter, file search, image gen, "
     "remote + local MCP, agent-as-tool, approvals.", FG),
    ("\u2022   4  multi-agent workflow patterns  \u2014  sequential, concurrent, handoff, group-chat.", FG),
    ("\u2022      hosting + OpenTelemetry  \u2014  Responses lifecycle, health, platform headers, "
     "graceful shutdown.", FG),
], size=16)
_, tf = box(s, 0.8, 6.7, 11.9, 0.7)
set_run(tf.paragraphs[0].add_run(),
        "One middleware seam here \u2014 the framework is a whole programming model.", 17, ACCENT,
        bold=True)

out = r"C:\Users\weidxu\OneDrive - Microsoft\doc8 - hackathon2026\demo.pptx"
prs.save(out)
print("saved", out, "slides:", len(prs.slides.__iter__.__self__._sldIdLst))
